import streamlit as st
import os
import shutil
from dotenv import load_dotenv
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
import pdfplumber
from langchain_huggingface import HuggingFaceEndpoint

# ---------------------- Environment ----------------------
load_dotenv()
token = os.getenv("HUGGINGFACEHUB_API_TOKEN")
UPLOAD_DIR = "uploaded_docs"
PROCESSED_DIR = "processed_docs"
os.makedirs(UPLOAD_DIR, exist_ok=True)
os.makedirs(PROCESSED_DIR, exist_ok=True)

# ---------------------- Mock User DB ----------------------
users_db = {
    "admin": {"password": "admin123", "role": "admin"},
    "user1": {"password": "user123", "role": "user"},
}

# ---------------------- State Init ----------------------
if "authenticated" not in st.session_state:
    st.session_state.authenticated = False
    st.session_state.username = ""
    st.session_state.role = ""

# ---------------------- Auth Logic ----------------------
def login(username, password):
    user = users_db.get(username)
    if user and user["password"] == password:
        st.session_state.authenticated = True
        st.session_state.username = username
        st.session_state.role = user["role"]
        return True
    return False

# ---------------------- LLM ----------------------
llm = HuggingFaceEndpoint(
    repo_id="mistralai/Mixtral-8x7B-Instruct-v0.1",
    task="text-generation",
    huggingfacehub_api_token=token,
    temperature=0.7,
    max_new_tokens=300,
)

# ---------------------- File Extraction ----------------------
def extract_text(file_path, file_type):
    try:
        if file_type == "pdf":
            with pdfplumber.open(file_path) as pdf:
                return "\n".join([p.extract_text() for p in pdf.pages if p.extract_text()])
        elif file_type == "docx":
            doc = Document(file_path)
            return "\n".join([p.text for p in doc.paragraphs])
        elif file_type == "txt":
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()
        elif file_type == "xlsx":
            wb = load_workbook(file_path)
            return "\n".join(
                " ".join([str(cell) if cell else "" for cell in row])
                for sheet in wb.worksheets
                for row in sheet.iter_rows(values_only=True)
            )
        elif file_type == "pptx":
            prs = Presentation(file_path)
            return "\n".join(
                shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")
            )
        else:
            return "Unsupported file type."
    except Exception as e:
        return f"Error extracting text: {str(e)}"

# ---------------------- File Utilities ----------------------
def safe_listdir(directory):
    try:
        return os.listdir(directory)
    except Exception as e:
        st.error(f"Error listing files in {directory}: {str(e)}")
        return []

def safe_read_file(filepath):
    try:
        with open(filepath, "r", encoding="utf-8") as f:
            return f.read()
    except Exception as e:
        st.error(f"Error reading file {filepath}: {str(e)}")
        return ""

def safe_remove_file(filepath):
    try:
        os.remove(filepath)
        return True
    except Exception as e:
        st.error(f"Error deleting file {filepath}: {str(e)}")
        return False

def save_uploaded_file(uploaded_file, save_path):
    try:
        with open(save_path, "wb") as f:
            shutil.copyfileobj(uploaded_file, f)
        return True
    except Exception as e:
        st.error(f"Error saving file {uploaded_file.name}: {str(e)}")
        return False

def save_processed_text(processed_path, text):
    try:
        with open(processed_path, "w", encoding="utf-8") as f:
            f.write(text)
        return True
    except Exception as e:
        st.error(f"Error saving processed text: {str(e)}")
        return False

# ---------------------- UI Function ----------------------
def run_ui():
    st.set_page_config(page_title="Document Search Bot", layout="wide")
    st.title("📄 Document Search Bot")

    # Login
    if not st.session_state.authenticated:
        st.sidebar.header("🔐 Login")
        username = st.sidebar.text_input("Username")
        password = st.sidebar.text_input("Password", type="password")
        if st.sidebar.button("Login"):
            if login(username, password):
                st.sidebar.success(f"Welcome {username} ({st.session_state.role})")
                st.rerun()
            else:
                st.sidebar.error("Invalid credentials.")
        st.stop()

    # Sidebar user info
    st.sidebar.success(f"Logged in as: {st.session_state.username} ({st.session_state.role})")
    nav = st.sidebar.radio("Navigation", ["Upload", "Query", "Manage Files"])

    # ---------------------- Upload (Admin only) ----------------------
    if nav == "Upload":
        st.header("📤 Upload Documents")
        if st.session_state.role != "admin":
            st.warning("Only admins can upload files.")
        else:
            file = st.file_uploader("Upload a document (PDF, DOCX, TXT, XLSX, PPTX)", type=["pdf", "docx", "txt", "xlsx", "pptx"])
            if file:
                if file.size > 2 * 1024 * 1024:
                    st.error("File exceeds 2MB limit.")
                else:
                    file_path = os.path.join(UPLOAD_DIR, file.name)
                    if save_uploaded_file(file, file_path):
                        ext = file.name.split(".")[-1].lower()
                        text = extract_text(file_path, ext)
                        processed_path = os.path.join(PROCESSED_DIR, file.name + ".txt")
                        if save_processed_text(processed_path, text):
                            st.success(f"{file.name} uploaded and processed successfully!")

    # ---------------------- Query ----------------------
    elif nav == "Query":
        st.header("🔍 Ask a Question")
        processed_files = safe_listdir(PROCESSED_DIR)
        if not processed_files:
            st.info("No documents available.")
        else:
            documents = {}
            for fname in processed_files:
                documents[fname] = safe_read_file(os.path.join(PROCESSED_DIR, fname))

            query = st.text_input("Your question:")
            if st.button("Search") and query:
                context = "\n\n".join(list(documents.values())[:3])
                prompt = f"""Use the following context to answer the question.
If the answer is not found, respond with 'Answer not found in the documents.'

Context:
{context}

Question: {query}
Answer:"""
                try:
                    result = llm.invoke(prompt)
                    st.markdown("**Answer:**")
                    st.write(result)
                except Exception as e:
                    st.error(f"Model error: {str(e)}")

    # ---------------------- Manage Files ----------------------
    elif nav == "Manage Files":
        st.header("📂 Uploaded Documents")
        processed_files = safe_listdir(PROCESSED_DIR)
        if not processed_files:
            st.info("No files uploaded.")
        else:
            for file in processed_files:
                col1, col2 = st.columns([4, 1])
                col1.write(file.replace(".txt", ""))
                if st.session_state.role == "admin":
                    if col2.button("❌", key=file):
                        if safe_remove_file(os.path.join(PROCESSED_DIR, file)):
                            st.success(f"Deleted {file.replace('.txt', '')}")
                            st.rerun()

# ---------------------- Main Function ----------------------
def main():
    run_ui()

if __name__ == "__main__":
    main()